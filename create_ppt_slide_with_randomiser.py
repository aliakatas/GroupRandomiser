import tkinter as tk
from tkinter import ttk, filedialog
import random
from pptx import Presentation

####################################
def append_to_powerpoint(pptx_file, txt, mod_pptx_file=None):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Create a new slide with a title and content layout
    slide_layout = presentation.slide_layouts[5]  # 5 corresponds to Title and Content layout
    new_slide = presentation.slides.add_slide(slide_layout)

    # Set the title and content for the new slide
    title = new_slide.shapes.title
    title.text = "Random Grouping!"

    for shape in new_slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.clear()  # remove any existing paragraphs, leaving one empty one
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = txt

    if mod_pptx_file is not None:
        # Save the modified presentation
        presentation.save(mod_pptx_file)
        result_table.insert("", "end", values=("Saved to", f"{mod_pptx_file}"))
    else:
        presentation.save(pptx_file)
        result_table.insert("", "end", values=("Saved to", f"{pptx_file}"))

####################################
def create_random_groups(n, group_size, allow_less=False):
    names = [f'{i + 1}' for i in range(n)]

    groups = list()
    if group_size > n:
        return groups

    while n >= group_size:
        group = list()
        for i in range(group_size):
            name = random.choice(names)
            group.append(name)
            names.remove(name)
        
        groups.append(group)
        n = len(names)
    
    if allow_less:
        groups.append(names)
    else:
        ngrp = len(groups)
        for i, name in enumerate(names):
            i = i % ngrp
            groups[i].append(name)
    
    return groups

# Function to calculate results based on input values
def calculate_results():
    # Get the input values as strings
    num_of_students_val = num_of_students.get()
    num_students_group_val = num_students_group.get()
    allow_less_in_group_val = allow_less_in_group.get()

    if not (num_of_students_val.isnumeric() and num_students_group_val.isnumeric()):
        result_table.delete(*result_table.get_children())  # Clear existing rows
        result_table.insert("", "end", values=("Please enter values for", "Number of students"))
        result_table.insert("", "end", values=("Please enter values for", "Number of students in group"))
        return

    # Convert input values to float (you can add error handling)
    num_of_students_val = int(num_of_students_val)
    num_students_group_val = int(num_students_group_val)
    allow_less_in_group_val = int(allow_less_in_group_val) > 0

    # Process inputs
    groups = create_random_groups(num_of_students_val, num_students_group_val, allow_less=allow_less_in_group_val)

    # Display the results in the table
    result_table.delete(*result_table.get_children())  # Clear existing rows
    global slide_txt
    
    slide_txt = ""
    for idx, group in enumerate(groups):
        result_table.insert("", "end", values=(f"Group {idx + 1} ({len(group)})", ", ".join(group)))
        slide_txt += f"Group {idx + 1} ({len(group)}) \t\t {', '.join(group)}\n"

####################################
def save():
    file_path = file_path_entry.get()

    if file_path:
        append_to_powerpoint(file_path, slide_txt)
        
    else:
        result_table.delete(*result_table.get_children())  # Clear existing rows
        result_table.insert("", "end", values=(f"Output file:", "Not selected!"))

# Function to open a file dialog and set the file path
def browse_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        file_path_entry.delete(0, tk.END)  # Clear the entry field
        file_path_entry.insert(0, file_path)  # Insert the selected file path

#############################################
if __name__ == '__main__':
    groups = list()
    slide_txt = ""

    # Create the main application window
    root = tk.Tk()
    root.title("Simple Randomiser")

    # Create labels and entry widgets for input values
    label1 = tk.Label(root, text="Number of students:")
    label1.pack()
    num_of_students = tk.Entry(root)
    num_of_students.pack()

    label2 = tk.Label(root, text="Number of students in group:")
    label2.pack()
    num_students_group = tk.Entry(root)
    num_students_group.pack()

    # Create a checkbox for on-off control
    allow_less_in_group = tk.IntVar()  # Variable to hold the checkbox state (0 or 1)
    checkbox = tk.Checkbutton(root, text="Allow less in group", variable=allow_less_in_group)
    checkbox.pack()

    # Create a button to trigger the calculation
    calculate_button = tk.Button(root, text="Calculate", command=calculate_results)
    calculate_button.pack()

    # Create a label and entry field for the file path
    file_path_label = tk.Label(root, text="File Path:")
    file_path_label.pack()

    file_path_entry = tk.Entry(root)
    file_path_entry.pack()

    # Create a button to open the file dialog
    browse_button = tk.Button(root, text="Browse", command=browse_file)
    browse_button.pack()

    # Create a button to trigger saving
    save_button = tk.Button(root, text="Save", command=save)
    save_button.pack()

    # Create a table to display the results
    columns = ("Description", "Result")
    result_table = ttk.Treeview(root, columns=columns, show="headings")
    result_table.heading("Description", text="Description")
    result_table.heading("Result", text="Result")
    result_table.pack()

    # Run the application
    root.mainloop()


